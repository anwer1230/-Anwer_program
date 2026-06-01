"""
محول HTML إلى PowerPoint
يقرأ كود HTML ويستخرج الشرائح والتنسيقات والألوان وينتج ملف PPTX
"""

import re
import tempfile
from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUTS_DIR = Path(__file__).parent.parent / "outputs"
OUTPUTS_DIR.mkdir(exist_ok=True)


# ── مساعدات اللون ─────────────────────────────────────────────────────────────

def _hex_to_rgb(hex_str: str) -> Tuple[int, int, int]:
    hex_str = hex_str.strip().lstrip("#")
    if len(hex_str) == 3:
        hex_str = "".join(c * 2 for c in hex_str)
    if len(hex_str) == 6:
        try:
            return (int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
        except Exception:
            pass
    return (102, 126, 234)


def _css_color_to_rgb(color_str: str) -> Optional[Tuple[int, int, int]]:
    if not color_str:
        return None
    color_str = color_str.strip()
    if color_str.startswith("#"):
        return _hex_to_rgb(color_str)
    m = re.match(r"rgb\(\s*(\d+),\s*(\d+),\s*(\d+)\)", color_str)
    if m:
        return (int(m.group(1)), int(m.group(2)), int(m.group(3)))
    named = {
        "white": (255, 255, 255), "black": (0, 0, 0),
        "red": (220, 50, 50), "blue": (66, 133, 244),
        "green": (52, 168, 83), "purple": (106, 27, 154),
        "orange": (255, 109, 0), "gray": (117, 117, 117),
        "navy": (26, 35, 126), "teal": (0, 121, 107),
        "darkblue": (13, 71, 161), "lightblue": (100, 181, 246),
    }
    return named.get(color_str.lower())


def _rgb(t: Tuple[int, int, int]) -> RGBColor:
    return RGBColor(*t)


# ── استخراج CSS من الـ HTML ───────────────────────────────────────────────────

def _extract_css_vars(html: str) -> Dict[str, str]:
    """استخراج متغيرات CSS ومتغيرات اللون من الـ HTML"""
    css_blocks = re.findall(r"<style[^>]*>(.*?)</style>", html, re.DOTALL | re.IGNORECASE)
    css_text = " ".join(css_blocks)
    vars_map: Dict[str, str] = {}
    for m in re.finditer(r"--([\w-]+)\s*:\s*([^;]+);", css_text):
        vars_map[f"--{m.group(1)}"] = m.group(2).strip()
    return vars_map


def _extract_inline_style(tag) -> Dict[str, str]:
    style_str = tag.get("style", "")
    result: Dict[str, str] = {}
    for part in style_str.split(";"):
        if ":" in part:
            k, v = part.split(":", 1)
            result[k.strip().lower()] = v.strip()
    return result


# ── تحليل HTML ────────────────────────────────────────────────────────────────

def parse_html_to_slides(html: str) -> List[Dict[str, Any]]:
    """
    يحوّل HTML إلى قائمة شرائح مع بيانات التنسيق.
    يدعم: section / .slide / [data-slide] / article / div.slide
    """
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(html, "lxml")

    # إزالة السكريبتات
    for tag in soup.find_all(["script", "noscript"]):
        tag.decompose()

    css_vars = _extract_css_vars(html)

    # البحث عن الشرائح بأولوية
    slide_candidates = (
        soup.find_all("section")
        or soup.find_all(class_=re.compile(r"slide|page|frame", re.I))
        or soup.find_all(attrs={"data-slide": True})
        or soup.find_all("article")
    )

    # لو ما في شرائح واضحة، نقسّم بناءً على H2/H3
    if not slide_candidates:
        slide_candidates = _split_by_headings(soup)

    slides: List[Dict[str, Any]] = []
    for el in slide_candidates:
        slide = _extract_slide_data(el, css_vars)
        if slide.get("title") or slide.get("bullets"):
            slides.append(slide)

    # fallback: كل المحتوى شريحة واحدة
    if not slides:
        slides = [_extract_slide_data(soup.body or soup, css_vars)]

    return slides


def _split_by_headings(soup) -> list:
    """تقسيم المحتوى إلى أقسام بناءً على العناوين"""
    from bs4 import BeautifulSoup, Tag
    sections = []
    current: Optional[list] = None
    for el in (soup.body or soup).children:
        if not hasattr(el, "name"):
            continue
        if el.name in ("h1", "h2", "h3"):
            if current:
                sections.append(current)
            current = [el]
        elif current is not None:
            current.append(el)
    if current:
        sections.append(current)
    # حوّل كل قائمة إلى عنصر بسيط
    result = []
    for group in sections:
        from bs4 import BeautifulSoup
        wrapper = BeautifulSoup("<div></div>", "lxml").div
        for el in group:
            wrapper.append(el.__copy__())
        result.append(wrapper)
    return result


def _extract_slide_data(el, css_vars: Dict[str, str]) -> Dict[str, Any]:
    """استخراج بيانات شريحة واحدة"""
    slide: Dict[str, Any] = {
        "title": "",
        "subtitle": "",
        "bullets": [],
        "slide_type": "bullets",
        "bg_color": None,
        "text_color": None,
        "accent_color": None,
        "table_data": None,
        "is_title_slide": False,
    }

    style = _extract_inline_style(el)

    # خلفية
    bg = style.get("background-color") or style.get("background")
    if bg:
        c = _css_color_to_rgb(bg.split()[0])
        if c:
            slide["bg_color"] = c

    # لون النص
    tc = style.get("color")
    if tc:
        c = _css_color_to_rgb(tc)
        if c:
            slide["text_color"] = c

    # العنوان
    for htag in ["h1", "h2", "h3", "h4"]:
        h = el.find(htag)
        if h:
            slide["title"] = h.get_text(strip=True)
            if htag == "h1":
                slide["is_title_slide"] = True
            break

    # العنوان الفرعي
    for stag in ["h2", "h3", ".subtitle", "p.subtitle"]:
        if "." in stag:
            cls = stag.split(".")[1]
            s = el.find(class_=cls)
        else:
            s = el.find(stag)
        if s and s.get_text(strip=True) != slide["title"]:
            slide["subtitle"] = s.get_text(strip=True)
            break

    # النقاط
    bullets = []
    for li in el.find_all("li"):
        txt = li.get_text(strip=True)
        if txt:
            bullets.append(txt)

    # فقرات (لو ما في قوائم)
    if not bullets:
        for p in el.find_all("p"):
            txt = p.get_text(strip=True)
            if txt and txt != slide["title"] and txt != slide["subtitle"] and len(txt) > 3:
                bullets.append(txt)

    slide["bullets"] = bullets[:6]

    # جداول
    tbl = el.find("table")
    if tbl:
        rows = []
        for tr in tbl.find_all("tr"):
            cells = [td.get_text(strip=True) for td in tr.find_all(["td", "th"])]
            if cells:
                rows.append(cells)
        if rows:
            slide["table_data"] = rows
            slide["slide_type"] = "table"

    # تحديد نوع الشريحة
    classes = " ".join(el.get("class", []))
    if re.search(r"title|cover|first|intro", classes, re.I) or slide["is_title_slide"]:
        slide["slide_type"] = "title"
    elif re.search(r"end|outro|thank|conclusion", classes, re.I):
        slide["slide_type"] = "conclusion"
    elif slide["table_data"]:
        slide["slide_type"] = "table"

    return slide


# ── بناء PPTX ────────────────────────────────────────────────────────────────

DEFAULT_THEME = {
    "primary": (102, 126, 234),
    "secondary": (118, 75, 162),
    "accent": (255, 255, 255),
    "text_dark": (26, 26, 46),
    "text_light": (255, 255, 255),
    "bg": (245, 247, 255),
}


def _pick_colors(slide_data: Dict, theme: Dict) -> Dict:
    """اختار ألوان الشريحة من البيانات المستخرجة أو القالب الافتراضي"""
    t = dict(theme)
    if slide_data.get("bg_color"):
        t["primary"] = slide_data["bg_color"]
        t["secondary"] = slide_data["bg_color"]
    if slide_data.get("text_color"):
        t["text_light"] = slide_data["text_color"]
        t["text_dark"] = slide_data["text_color"]
    return t


def html_to_pptx(html: str, override_title: str = "") -> str:
    slides_data = parse_html_to_slides(html)
    if not slides_data:
        raise ValueError("لم يتم العثور على محتوى قابل للتحويل في الكود HTML المدخل")

    if override_title and slides_data:
        slides_data[0]["title"] = override_title

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for i, sd in enumerate(slides_data):
        t = _pick_colors(sd, DEFAULT_THEME)
        stype = sd.get("slide_type", "bullets")
        if stype == "title":
            _build_title_slide(prs, sd, t)
        elif stype == "table" and sd.get("table_data"):
            _build_table_slide(prs, sd, t)
        elif stype == "conclusion":
            _build_conclusion_slide(prs, sd, t)
        else:
            _build_content_slide(prs, sd, t, i)

    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out = str(OUTPUTS_DIR / f"html_presentation_{ts}.pptx")
    prs.save(out)
    return out


# ── بناء الشرائح ─────────────────────────────────────────────────────────────

def _txt(slide, l, t, w, h, text, size, bold=False, color=(255,255,255), align=PP_ALIGN.RIGHT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = _rgb(color)
    return box


def _rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = _rgb(color)
    s.line.fill.background()
    return s


def _build_title_slide(prs, sd, t):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background.fill; bg.solid(); bg.fore_color.rgb = _rgb(t["primary"])
    _rect(sl, 0, 6.7, 13.33, 0.8, t["secondary"])
    _txt(sl, 1, 2.0, 11.3, 1.9, sd["title"] or "العنوان", 44, bold=True, color=t["text_light"], align=PP_ALIGN.CENTER)
    if sd.get("subtitle"):
        _txt(sl, 1.5, 4.1, 10.3, 0.9, sd["subtitle"], 24, color=(210,210,240), align=PP_ALIGN.CENTER)
    line = sl.shapes.add_shape(1, Inches(4.2), Inches(4.9), Inches(4.9), Emu(38000))
    line.fill.solid(); line.fill.fore_color.rgb = _rgb(t["accent"]); line.line.fill.background()


def _build_content_slide(prs, sd, t, idx):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background.fill; bg.solid(); bg.fore_color.rgb = _rgb(t["bg"])
    _rect(sl, 0, 0, 13.33, 1.35, t["primary"])
    _txt(sl, 0.4, 0.15, 12, 1.05, sd["title"] or f"شريحة {idx+1}", 27, bold=True, color=t["text_light"], align=PP_ALIGN.RIGHT)

    bullets = sd.get("bullets", [])
    box = sl.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.6))
    tf = box.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT; p.space_before = Pt(9)
        r = p.add_run(); r.text = f"◆  {b}"
        r.font.size = Pt(20); r.font.color.rgb = _rgb(t["text_dark"])

    # رقم
    nb = sl.shapes.add_textbox(Inches(12.6), Inches(6.9), Inches(0.7), Inches(0.5))
    np_ = nb.text_frame.paragraphs[0]; np_.alignment = PP_ALIGN.CENTER
    nr = np_.add_run(); nr.text = str(idx + 1)
    nr.font.size = Pt(11); nr.font.color.rgb = _rgb(t["secondary"])


def _build_table_slide(prs, sd, t):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background.fill; bg.solid(); bg.fore_color.rgb = _rgb(t["bg"])
    _rect(sl, 0, 0, 13.33, 1.35, t["primary"])
    _txt(sl, 0.4, 0.15, 12, 1.05, sd["title"] or "جدول البيانات", 27, bold=True, color=t["text_light"], align=PP_ALIGN.RIGHT)

    td = sd["table_data"]
    rows, cols = len(td), max(len(r) for r in td)
    if rows == 0 or cols == 0:
        return

    tbl = sl.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.6),
        Inches(12.3), Inches(min(5.5, 0.5 + rows * 0.55))
    ).table

    for ri, row in enumerate(td):
        for ci in range(cols):
            cell = tbl.cell(ri, ci)
            cell.text = row[ci] if ci < len(row) else ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            if p.runs:
                p.runs[0].font.size = Pt(14)
                p.runs[0].font.bold = (ri == 0)
                p.runs[0].font.color.rgb = _rgb(t["text_light"] if ri == 0 else t["text_dark"])
            _set_cell_bg(cell, t["primary"] if ri == 0 else ((235,238,255) if ri%2==0 else (255,255,255)))


def _build_conclusion_slide(prs, sd, t):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background.fill; bg.solid(); bg.fore_color.rgb = _rgb(t["secondary"])
    _txt(sl, 1, 1.5, 11.3, 1.2, sd["title"] or "الخلاصة", 36, bold=True, color=t["text_light"], align=PP_ALIGN.CENTER)
    line = sl.shapes.add_shape(1, Inches(3.5), Inches(2.9), Inches(6.3), Emu(35000))
    line.fill.solid(); line.fill.fore_color.rgb = _rgb(t["accent"]); line.line.fill.background()
    box = sl.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(10.3), Inches(4.0))
    tf = box.text_frame; tf.word_wrap = True
    for i, b in enumerate(sd.get("bullets", [])):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT; p.space_before = Pt(12)
        r = p.add_run(); r.text = b
        r.font.size = Pt(22); r.font.color.rgb = RGBColor(220, 220, 255)


def _set_cell_bg(cell, color_tuple):
    from pptx.oxml.ns import qn
    from lxml import etree
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    sf = etree.SubElement(tcPr, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", "{:02X}{:02X}{:02X}".format(*color_tuple))


from pptx.dml.color import RGBColor
