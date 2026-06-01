"""
مولد العروض التقديمية - ينشئ ملفات PowerPoint احترافية
يدعم: صفحة الغلاف، الجداول، الرسوم البيانية، RTL عربي صحيح
"""

import os
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Any, Tuple, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

from modules.templates import ProfessionalTemplates, THEMES

OUTPUTS_DIR = Path(__file__).parent.parent / "outputs"
OUTPUTS_DIR.mkdir(exist_ok=True)

# الخطوط المدعومة مع العربية
ARABIC_FONTS = {
    "Traditional Arabic": "Traditional Arabic",
    "Simplified Arabic":  "Simplified Arabic",
    "Arial":              "Arial",
    "Dubai":              "Dubai",
    "Calibri":            "Calibri",
    "Times New Roman":    "Times New Roman",
}


def _rgb(t: Tuple[int, int, int]) -> RGBColor:
    return RGBColor(*t)


def _set_rtl(paragraph):
    """تفعيل اتجاه RTL على الفقرة عبر XML مباشرةً"""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")


def _set_run_font(run, font_name: str):
    """تعيين الخط للنص العادي والنص العربي (complex script)"""
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    # خط اللاتيني
    latin = rPr.find(qn("a:latin"))
    if latin is None:
        latin = etree.SubElement(rPr, qn("a:latin"))
    latin.set("typeface", font_name)
    # خط العربي (complex script) — الأهم لعرض الحروف صحيحاً
    cs = rPr.find(qn("a:cs"))
    if cs is None:
        cs = etree.SubElement(rPr, qn("a:cs"))
    cs.set("typeface", font_name)
    # خط الشرق الأقصى
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font_name)


def _add_text_box(slide, left, top, width, height,
                  text, size, bold=False, color=(0, 0, 0),
                  align=PP_ALIGN.RIGHT, wrap=True,
                  font_name="Traditional Arabic"):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    _set_rtl(p)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    _set_run_font(run, font_name)
    return box


class PresentationGenerator:
    def __init__(self):
        self.templates = ProfessionalTemplates()

    def create_presentation(
        self,
        slides_data: List[Dict[str, Any]],
        theme_color: str = "blue",
        cover_data: Optional[Dict[str, str]] = None,
        extracted_images: Optional[List[str]] = None,
        font_name: str = "Traditional Arabic",
        body_font_size: int = 22,
        ai_images: bool = False,
        groq_client=None,
    ) -> str:
        self.font_name      = font_name
        self.body_font_size = body_font_size

        theme = self.templates.get_theme(theme_color)
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        if cover_data:
            self._add_cover_slide(prs, cover_data, theme)

        for i, slide_data in enumerate(slides_data):
            stype = slide_data.get("slide_type", "bullets")
            if stype == "title":
                self._add_title_slide(prs, slide_data, theme)
            elif stype == "table":
                self._add_table_slide(prs, slide_data, theme)
            elif stype == "chart":
                self._add_chart_slide(prs, slide_data, theme)
            elif stype == "conclusion":
                self._add_conclusion_slide(prs, slide_data, theme)
            else:
                # ── تحديد مصدر الصورة ──
                img_path = None
                if ai_images:
                    try:
                        from modules.image_fetcher import get_slide_image
                        img_path = get_slide_image(
                            slide_title  = slide_data.get("title", ""),
                            slide_bullets= slide_data.get("bullets", []),
                            slide_type   = stype,
                            groq_client  = groq_client,
                        )
                    except Exception:
                        img_path = None
                elif extracted_images and i < len(extracted_images):
                    img_path = extracted_images[i]

                self._add_content_slide(prs, slide_data, theme, i, img_path)

        timestamp   = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = str(OUTPUTS_DIR / f"عرض_{timestamp}.pptx")
        prs.save(output_path)
        return output_path

    # ── صفحة الغلاف ──────────────────────────────────────────
    def _add_cover_slide(self, prs, cover_data, theme):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _rgb(theme["primary"])

        rect = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(2.2))
        rect.fill.solid()
        rect.fill.fore_color.rgb = _rgb(theme["secondary"])
        rect.line.fill.background()

        logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(2), Inches(1.5))
        lp = logo_box.text_frame.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lr = lp.add_run()
        lr.text = cover_data.get("logo", "📊")
        lr.font.size = Pt(52)

        org = cover_data.get("organization", "")
        if org:
            org_box = slide.shapes.add_textbox(Inches(2.8), Inches(0.5), Inches(10), Inches(1))
            op = org_box.text_frame.paragraphs[0]
            op.alignment = PP_ALIGN.RIGHT
            _set_rtl(op)
            or_ = op.add_run()
            or_.text = org
            or_.font.size = Pt(20)
            or_.font.bold = True
            or_.font.color.rgb = RGBColor(220, 220, 255)
            _set_run_font(or_, self.font_name)

        fn = self.font_name
        _add_text_box(slide, 1, 2.5, 11.3, 1.8,
                      cover_data.get("title", "العنوان الرئيسي"),
                      42, bold=True, color=theme["text_light"],
                      align=PP_ALIGN.CENTER, font_name=fn)

        subtitle = cover_data.get("subtitle", "")
        if subtitle:
            _add_text_box(slide, 1.5, 4.4, 10.3, 0.9,
                          subtitle, 22, color=(210, 210, 240),
                          align=PP_ALIGN.CENTER, font_name=fn)

        line = slide.shapes.add_shape(1, Inches(3.5), Inches(5.5), Inches(6.3), Emu(35000))
        line.fill.solid()
        line.fill.fore_color.rgb = _rgb(theme["accent"])
        line.line.fill.background()

        meta_parts = []
        if cover_data.get("presenter"):
            meta_parts.append(f"إعداد: {cover_data['presenter']}")
        if cover_data.get("date"):
            meta_parts.append(cover_data["date"])
        if meta_parts:
            _add_text_box(slide, 1, 5.9, 11.3, 0.7,
                          "  |  ".join(meta_parts), 16,
                          color=(200, 200, 230), align=PP_ALIGN.CENTER, font_name=fn)

        note = cover_data.get("note", "")
        if note:
            _add_text_box(slide, 1.5, 6.6, 10.3, 0.7,
                          note, 14, color=(180, 180, 220),
                          align=PP_ALIGN.CENTER, font_name=fn)

    # ── شريحة العنوان ────────────────────────────────────────
    def _add_title_slide(self, prs, data, theme):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _rgb(theme["primary"])

        self._add_bottom_strip(slide, theme)
        fn = self.font_name

        _add_text_box(slide, 1, 2.2, 11.3, 1.8,
                      data.get("title", "العنوان"),
                      44, bold=True, color=theme["text_light"],
                      align=PP_ALIGN.CENTER, font_name=fn)

        subtitle = data.get("subtitle", "")
        if subtitle:
            _add_text_box(slide, 2, 4.2, 9.3, 1,
                          subtitle, 24, color=(220, 220, 255),
                          align=PP_ALIGN.CENTER, font_name=fn)

        line = slide.shapes.add_shape(1, Inches(4.5), Inches(4.0), Inches(4.3), Emu(40000))
        line.fill.solid()
        line.fill.fore_color.rgb = _rgb(theme["accent"])
        line.line.fill.background()

    # ── شريحة المحتوى (نقاط) ────────────────────────────────
    def _add_content_slide(self, prs, data, theme, index, img_path=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _rgb(theme["bg"])

        # ── رأس ──
        header = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.33), Inches(1.35)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = _rgb(theme["primary"])
        header.line.fill.background()

        _add_text_box(slide, 0.4, 0.15, 12, 1.05,
                      data.get("title", ""), 27, bold=True,
                      color=theme["text_light"], align=PP_ALIGN.RIGHT,
                      font_name=self.font_name)

        bullets = data.get("bullets", [])

        # ── تخطيط مع صورة على اليسار ──
        has_image = img_path and Path(img_path).exists()
        if has_image:
            # نصف أيمن للنص (RTL: النص على اليمين)
            text_left  = 4.8
            text_width = 8.1
            content_box = slide.shapes.add_textbox(
                Inches(text_left), Inches(1.55), Inches(text_width), Inches(5.7)
            )
            self._fill_bullets(content_box.text_frame, bullets, theme)

            # الصورة على اليسار مع إطار ملوّن
            try:
                img_left = 0.25
                img_top  = 1.55
                img_w    = 4.3
                img_h    = 5.5
                # إطار خلفي
                frame = slide.shapes.add_shape(
                    1, Inches(img_left - 0.07), Inches(img_top - 0.07),
                    Inches(img_w + 0.14), Inches(img_h + 0.14)
                )
                frame.fill.solid()
                frame.fill.fore_color.rgb = _rgb(theme["accent"])
                frame.line.fill.background()

                slide.shapes.add_picture(
                    img_path,
                    Inches(img_left), Inches(img_top),
                    Inches(img_w), Inches(img_h)
                )
            except Exception:
                # عند فشل الصورة: نص يملأ الشريحة
                content_box2 = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.55), Inches(12.3), Inches(5.7)
                )
                self._fill_bullets(content_box2.text_frame, bullets, theme)
        else:
            # بدون صورة: النص يملأ الشريحة كاملة
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.55), Inches(12.3), Inches(5.7)
            )
            self._fill_bullets(content_box.text_frame, bullets, theme)

        # رقم الشريحة
        num_box = slide.shapes.add_textbox(Inches(12.6), Inches(6.9), Inches(0.7), Inches(0.5))
        np_ = num_box.text_frame.paragraphs[0]
        np_.alignment = PP_ALIGN.CENTER
        nr = np_.add_run()
        nr.text = str(index + 1)
        nr.font.size = Pt(11)
        nr.font.color.rgb = _rgb(theme["secondary"])

    def _fill_bullets(self, tf, bullets, theme):
        """ملء النقاط مع RTL صحيح ومحاذاة يمين"""
        tf.word_wrap = True
        fs = self.body_font_size
        fn = self.font_name

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.RIGHT
            _set_rtl(p)                          # ← إصلاح RTL
            p.space_before = Pt(10)
            p.space_after  = Pt(4)

            # الرمز يُكتب أولاً ثم المسافة ثم النص — في RTL سيظهر الرمز على اليمين
            run = p.add_run()
            run.text = f"◆ {bullet}"
            run.font.size = Pt(fs)
            run.font.color.rgb = _rgb(theme["text_dark"])
            _set_run_font(run, fn)

    # ── شريحة الجدول ─────────────────────────────────────────
    def _add_table_slide(self, prs, data, theme):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _rgb(theme["bg"])

        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.35))
        header.fill.solid()
        header.fill.fore_color.rgb = _rgb(theme["primary"])
        header.line.fill.background()

        _add_text_box(slide, 0.4, 0.15, 12, 1.05,
                      data.get("title", "جدول البيانات"), 27, bold=True,
                      color=theme["text_light"], align=PP_ALIGN.RIGHT,
                      font_name=self.font_name)

        table_data: List[List[str]] = data.get("table_data", [])
        if not table_data:
            return

        rows = len(table_data)
        cols = max(len(r) for r in table_data)
        if rows == 0 or cols == 0:
            return

        tbl_left   = Inches(0.5)
        tbl_top    = Inches(1.6)
        tbl_width  = Inches(12.3)
        tbl_height = Inches(min(5.5, 0.5 + rows * 0.55))

        table = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_width, tbl_height).table

        fn = self.font_name
        for r_idx, row in enumerate(table_data):
            for c_idx in range(cols):
                cell      = table.cell(r_idx, c_idx)
                cell_text = row[c_idx] if c_idx < len(row) else ""
                cell.text = str(cell_text)

                tf  = cell.text_frame
                par = tf.paragraphs[0]
                par.alignment = PP_ALIGN.CENTER
                _set_rtl(par)

                run = par.runs
                if run:
                    run[0].font.size  = Pt(15)
                    run[0].font.bold  = (r_idx == 0)
                    run[0].font.color.rgb = (
                        _rgb(theme["text_light"]) if r_idx == 0 else _rgb(theme["text_dark"])
                    )
                    _set_run_font(run[0], fn)

                if r_idx == 0:
                    self._set_cell_bg(cell, theme["primary"])
                elif r_idx % 2 == 0:
                    self._set_cell_bg(cell, (235, 238, 255))
                else:
                    self._set_cell_bg(cell, (255, 255, 255))

    def _set_cell_bg(self, cell, color_tuple):
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        sf   = etree.SubElement(tcPr, qn("a:solidFill"))
        sc   = etree.SubElement(sf, qn("a:srgbClr"))
        sc.set("val", "{:02X}{:02X}{:02X}".format(*color_tuple))

    # backward-compat alias
    def _set_cell_background(self, cell, color_tuple):
        self._set_cell_bg(cell, color_tuple)

    # ── شريحة الرسم البياني ──────────────────────────────────
    def _add_chart_slide(self, prs, data, theme):
        from modules.chart_generator import create_chart_from_data

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _rgb(theme["bg"])

        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.35))
        header.fill.solid()
        header.fill.fore_color.rgb = _rgb(theme["primary"])
        header.line.fill.background()

        _add_text_box(slide, 0.4, 0.15, 12, 1.05,
                      data.get("title", "رسم بياني"), 27, bold=True,
                      color=theme["text_light"], align=PP_ALIGN.RIGHT,
                      font_name=self.font_name)

        chart_type  = data.get("chart_type", "bar")
        labels      = data.get("chart_labels", [])
        values      = data.get("chart_values", [])
        chart_title = data.get("chart_title", "")

        if not labels or not values:
            _add_text_box(slide, 1, 3, 11, 2,
                          "لا توجد بيانات كافية لرسم الرسم البياني",
                          20, color=theme["text_dark"], font_name=self.font_name)
            return

        try:
            values_float = [float(v) for v in values]
        except Exception:
            values_float = [1.0] * len(labels)

        primary_hex = "#{:02x}{:02x}{:02x}".format(*theme["primary"])
        img_path    = create_chart_from_data(chart_type, labels, values_float,
                                             chart_title, primary_hex)

        if img_path and Path(img_path).exists():
            slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.6),
                                     Inches(10.3), Inches(5.6))
        else:
            _add_text_box(slide, 1, 3, 11, 2,
                          "تعذّر إنشاء الرسم البياني", 20,
                          color=theme["text_dark"], font_name=self.font_name)

    # ── شريحة الخلاصة ────────────────────────────────────────
    def _add_conclusion_slide(self, prs, data, theme):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _rgb(theme["secondary"])

        fn = self.font_name
        _add_text_box(slide, 1, 1.5, 11.3, 1.2,
                      data.get("title", "الخلاصة"), 36, bold=True,
                      color=theme["text_light"], align=PP_ALIGN.CENTER,
                      font_name=fn)

        line = slide.shapes.add_shape(1, Inches(3.5), Inches(2.9), Inches(6.3), Emu(35000))
        line.fill.solid()
        line.fill.fore_color.rgb = _rgb(theme["accent"])
        line.line.fill.background()

        bullets   = data.get("bullets", [])
        cont_box  = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(10.3), Inches(4.0))
        tf        = cont_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.RIGHT
            _set_rtl(p)
            p.space_before = Pt(12)
            run = p.add_run()
            run.text = bullet
            run.font.size  = Pt(self.body_font_size)
            run.font.color.rgb = RGBColor(220, 220, 255)
            _set_run_font(run, fn)

    # ── مساعدات ──────────────────────────────────────────────
    def _add_bottom_strip(self, slide, theme):
        rect = slide.shapes.add_shape(1, Inches(0), Inches(6.8), Inches(13.33), Inches(0.7))
        rect.fill.solid()
        rect.fill.fore_color.rgb = _rgb(theme["secondary"])
        rect.line.fill.background()
